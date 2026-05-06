import os
import re
import logging
import uuid
from typing import List, Tuple, Optional
from pathlib import Path

from sqlalchemy.ext.asyncio import AsyncSession
from sqlalchemy import select, func

from app.models.documents import Document, DocumentChunk
from app.services.embedding_service import embed_batch
from app.core.config import settings

logger = logging.getLogger(__name__)


def extract_text_from_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    pages = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            pages.append(text)
    return "\n\n".join(pages)


def extract_text_from_docx(file_path: str) -> str:
    from docx import Document as DocxDocument
    doc = DocxDocument(file_path)
    return "\n\n".join(p.text for p in doc.paragraphs if p.text.strip())


def extract_text_from_pptx(file_path: str) -> str:
    import zipfile

    # First try python-pptx
    try:
        from pptx import Presentation
        prs = Presentation(file_path)
        slide_texts = []
        for slide in prs.slides:
            parts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    parts.append(shape.text.strip())
            if parts:
                slide_texts.append("\n".join(parts))
        return "\n\n".join(slide_texts)
    except Exception:
        pass

    # Fallback: extract text directly from XML inside the ZIP
    try:
        from xml.etree import ElementTree
        texts = []
        with zipfile.ZipFile(file_path) as zf:
            for name in sorted(zf.namelist()):
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    xml_data = zf.read(name)
                    root = ElementTree.fromstring(xml_data)
                    for elem in root.iter():
                        if elem.text and elem.text.strip():
                            texts.append(elem.text.strip())
        return "\n\n".join(texts)
    except Exception:
        pass

    raise ValueError(f"Could not extract text from PPTX: file may be corrupted")


def extract_text(file_path: str, file_type: str) -> str:
    extractors = {
        "pdf": extract_text_from_pdf,
        "docx": extract_text_from_docx,
        "pptx": extract_text_from_pptx,
        "txt": lambda p: Path(p).read_text(encoding="utf-8"),
    }
    extractor = extractors.get(file_type)
    if not extractor:
        raise ValueError(f"Unsupported file type: {file_type}")
    return extractor(file_path)


def chunk_text(text: str, chunk_size: int = 500, overlap: int = 50) -> List[Tuple[str, int]]:
    text = re.sub(r"\n{3,}", "\n\n", text).strip()
    if not text:
        return []

    paragraphs = re.split(r"\n\n+", text)
    chars_per_chunk = chunk_size * 4
    chars_overlap = overlap * 4

    chunks = []
    current = []
    current_len = 0

    for para in paragraphs:
        para = para.strip()
        if not para:
            continue

        para_len = len(para)

        if current_len + para_len > chars_per_chunk and current:
            chunk_str = "\n\n".join(current)
            chunks.append((chunk_str, len(chunk_str) // 4))

            overlap_text = chunk_str[-chars_overlap:] if chars_overlap > 0 else ""
            current = [overlap_text] if overlap_text else []
            current_len = len(overlap_text)

        current.append(para)
        current_len += para_len + 2

    if current:
        chunk_str = "\n\n".join(current)
        if chunk_str.strip():
            chunks.append((chunk_str, len(chunk_str) // 4))

    return chunks


def _get_upload_dir() -> Path:
    upload_dir = Path(settings.upload_dir).resolve()
    upload_dir.mkdir(parents=True, exist_ok=True)
    return upload_dir


def save_upload(content: bytes, extension: str) -> str:
    upload_dir = _get_upload_dir()
    filename = f"{uuid.uuid4()}.{extension}"
    file_path = str(upload_dir / filename)
    with open(file_path, "wb") as f:
        f.write(content)
    return file_path


def save_text_content(text: str) -> str:
    upload_dir = _get_upload_dir()
    filename = f"{uuid.uuid4()}.txt"
    file_path = str(upload_dir / filename)
    with open(file_path, "w", encoding="utf-8") as f:
        f.write(text)
    return file_path


async def create_document_record(
    db: AsyncSession,
    title: str,
    key_stage: str,
    subject: str,
    source_type: str,
    uploaded_by: int,
    exam_board: str = "None",
    tier: str = "None",
    unit_name: Optional[str] = None,
    source_url: Optional[str] = None,
    file_path: Optional[str] = None,
    file_type: Optional[str] = None,
    kb_type: str = "course_material",
) -> Document:
    doc = Document(
        title=title, key_stage=key_stage, subject=subject,
        exam_board=exam_board, tier=tier, unit_name=unit_name,
        source_type=source_type, uploaded_by=uploaded_by,
        source_url=source_url, file_path=file_path, file_type=file_type,
        kb_type=kb_type,
    )
    db.add(doc)
    await db.flush()
    await db.refresh(doc)
    return doc


async def process_document(db: AsyncSession, document_id: int) -> None:
    result = await db.execute(select(Document).where(Document.id == document_id))
    doc = result.scalar_one_or_none()
    if not doc:
        logger.error(f"Document {document_id} not found")
        return

    try:
        doc.status = "processing"
        await db.flush()

        file_type = doc.file_type or "txt"
        raw_text = extract_text(doc.file_path, file_type)

        if not raw_text.strip():
            raise ValueError("No text extracted from document")

        chunks = chunk_text(raw_text, settings.rag_chunk_size, settings.rag_chunk_overlap)

        if not chunks:
            raise ValueError("Document produced no text chunks")

        texts = [c[0] for c in chunks]
        embeddings = await embed_batch(texts)

        for idx, ((chunk_content, token_count), embedding) in enumerate(zip(chunks, embeddings)):
            chunk = DocumentChunk(
                document_id=doc.id,
                chunk_index=idx,
                content=chunk_content,
                embedding=embedding,
                token_count=token_count,
            )
            db.add(chunk)

        doc.chunk_count = len(chunks)
        doc.status = "ready"
        doc.error_message = None
        await db.flush()
        logger.info(f"Document {document_id} processed: {len(chunks)} chunks")

    except Exception as e:
        logger.error(f"Document {document_id} processing failed: {e}")
        doc.status = "failed"
        doc.error_message = str(e)[:500]
        await db.flush()


async def list_documents(
    db: AsyncSession,
    key_stage: Optional[str] = None,
    subject: Optional[str] = None,
    exam_board: Optional[str] = None,
    kb_type: Optional[str] = None,
    status: Optional[str] = None,
    limit: int = 100,
    offset: int = 0,
) -> Tuple[List[Document], int]:
    query = select(Document)
    count_query = select(func.count(Document.id))

    filters = []
    if key_stage:
        filters.append(Document.key_stage == key_stage)
    if subject:
        filters.append(Document.subject == subject)
    if exam_board:
        filters.append(Document.exam_board == exam_board)
    if kb_type:
        filters.append(Document.kb_type == kb_type)
    if status:
        filters.append(Document.status == status)

    for f in filters:
        query = query.where(f)
        count_query = count_query.where(f)

    query = query.order_by(Document.created_at.desc()).limit(limit).offset(offset)
    result = await db.execute(query)
    count_result = await db.execute(count_query)
    return list(result.scalars().all()), count_result.scalar() or 0


async def delete_document(db: AsyncSession, document_id: int) -> bool:
    result = await db.execute(select(Document).where(Document.id == document_id))
    doc = result.scalar_one_or_none()
    if not doc:
        return False

    if doc.file_path and os.path.exists(doc.file_path):
        try:
            os.unlink(doc.file_path)
        except OSError:
            pass

    await db.delete(doc)
    await db.flush()
    return True
